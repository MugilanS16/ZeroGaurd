import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Define some helper functions to add slides easily
def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # 0 is title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_content_slide(prs, title_text, content_points):
    slide_layout = prs.slide_layouts[1] # 1 is title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, point in enumerate(content_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            
# --- Slide Content ---

# Slide 1
add_title_slide(prs, "CyberCrimeAI Platform", "Modern Cybercrime Reporting System\n(25-Page Comprehensive Overview)")

# Slide 2
add_content_slide(prs, "1. Introduction to Cyber Crime", [
    "Cybercrimes are rapidly increasing globally.",
    "Traditional reporting methods are often slow and complex.",
    "Citizens face hurdles in understanding legal terminologies.",
    "A modernized, digital-first approach is necessary."
])

# Slide 3
add_content_slide(prs, "2. Project Vision & Mission", [
    "Vision: A safer digital space for every citizen.",
    "Mission: Simplify the reporting process using AI and modern tech.",
    "Goal: Provide immediate assistance, guidance, and formal reporting.",
    "Reduce the barrier to entry for cybercrime reporting."
])

# Slide 4
add_content_slide(prs, "3. Problem Statement", [
    "Victims of cybercrime often panic and don't know the exact steps to take.",
    "Law enforcement receives unstructured and incomplete reports.",
    "Lack of immediate automated guidance during an incident.",
    "Authentication and verification of reporters is a challenge."
])

# Slide 5
add_content_slide(prs, "4. Our Solution: CyberCrimeAI", [
    "An AI-powered web application for structured reporting.",
    "Interactive UI that guides the user step-by-step.",
    "Integrates SSO and DigiLocker for trusted identity verification.",
    "Automatically compiles evidence into a standardized PDF report."
])

# Slide 6
add_content_slide(prs, "5. Target Audience & Impact", [
    "General Citizens: Easy-to-use platform with no technical jargon.",
    "Law Enforcement: Receives organized, verifiable data.",
    "Impact: Faster response times and better tracking of incidents.",
    "Empowers individuals to take immediate action against fraud."
])

# Slide 7
add_content_slide(prs, "6. Entire Project Features - Overview", [
    "Modern UI/UX with responsive design.",
    "15 Specialized Crime Categories.",
    "Interactive AI Chat Sandbox.",
    "Automated Report Generation.",
    "Emergency Portal Integration.",
    "Secure Authentication flows."
])

# Slide 8
add_content_slide(prs, "7. Feature: Modern UI & Theme", [
    "Deep Blue primary colors for trust and security.",
    "Clean, white cards and backgrounds for readability.",
    "Subtle micro-animations to enhance user experience.",
    "Fully responsive for both desktop and mobile users."
])

# Slide 9
add_content_slide(prs, "8. Feature: 15 Specialized Categories", [
    "Platform covers 15 distinct types of cybercrimes.",
    "Includes Financial Fraud, Identity Theft, Cyberbullying, etc.",
    "Grid layout provides an easy overview of all options.",
    "Allows targeted data collection based on the crime type."
])

# Slide 10
add_content_slide(prs, "9. Feature: Live Search Capabilities", [
    "Users can instantly search for their specific issue.",
    "Search dynamically filters the 15 categories in real-time.",
    "Reduces time spent looking for the correct reporting section.",
    "Enhances accessibility for stressed users."
])

# Slide 11
add_content_slide(prs, "10. Feature: Interactive AI Chat Sandbox", [
    "A real-time AI assistant that acts as a first responder.",
    "Analyzes user input to understand the incident context.",
    "Provides immediate, actionable checklists (e.g., blocking cards).",
    "Acts as an empathetic guide during a stressful situation."
])

# Slide 12
add_content_slide(prs, "11. Feature: Dynamic Risk Leveling", [
    "AI evaluates the severity of the incident.",
    "Assigns risk levels (High, Medium, Low) dynamically.",
    "Helps prioritize severe incidents for immediate action.",
    "Visually alerts the user to critical next steps."
])

# Slide 13
add_content_slide(prs, "12. Feature: Dynamic Questioning", [
    "Static forms are replaced with AI-driven questions.",
    "Questions adapt based on the specific type of crime selected.",
    "Ensures only relevant information is requested from the user.",
    "Improves the quality and completeness of the final report."
])

# Slide 14
add_content_slide(prs, "13. Feature: Automated PDF Report", [
    "All user answers and uploaded evidence are compiled automatically.",
    "Generates a formal, legally structured PDF document.",
    "Includes timestamps, reference numbers, and user details.",
    "Ready for immediate submission to law enforcement."
])

# Slide 15
add_content_slide(prs, "14. Feature: Emergency Numbers Portal", [
    "One-tap access to critical helplines.",
    "Includes National Cyber Crime Helpline (1930).",
    "Includes General Emergency (112) and Women Helpline (1091).",
    "Crucial for immediate intervention in ongoing frauds."
])

# Slide 16
add_content_slide(prs, "15. Feature: Interactive Stats Counter", [
    "Displays real-time statistics of resolved cases or site usage.",
    "Numbers animate and count up automatically on scroll.",
    "Builds trust and credibility with new users.",
    "Demonstrates the platform's active usage."
])

# Slide 17
add_content_slide(prs, "16. Total Workflow - Overview", [
    "The application follows a structured 5-step process.",
    "Ensures users don't miss critical steps.",
    "Flow: Login -> Select Category -> AI Analysis -> Details -> Submit.",
    "Designed for maximum user retention and completion rates."
])

# Slide 18
add_content_slide(prs, "17. Workflow Step 1: User Initiation", [
    "User lands on the secure dashboard.",
    "Dashboard displays past reports and their current status.",
    "User clicks 'Report New Incident' to begin.",
    "System generates a unique tracking reference (CC-YYYY-NNNNN)."
])

# Slide 19
add_content_slide(prs, "18. Workflow Step 2: Category Selection", [
    "User is presented with the 15 crime categories.",
    "User selects the category that best matches their situation.",
    "This selection dictates the rest of the AI flow.",
    "Ensures specialized handling of the complaint."
])

# Slide 20
add_content_slide(prs, "19. Workflow Step 3: AI Incident Analysis", [
    "User interacts with the AI Chat Sandbox.",
    "Describes the incident in their own words.",
    "AI extracts key facts and suggests immediate preventative actions.",
    "Sets the stage for formal data collection."
])

# Slide 21
add_content_slide(prs, "20. Workflow Step 4: Evidence Submission", [
    "User answers the dynamic questions generated by the AI.",
    "Uploads relevant evidence (screenshots, bank statements, PDFs).",
    "Files are securely processed and attached to the session.",
    "System validates the completeness of the data."
])

# Slide 22
add_content_slide(prs, "21. Workflow Step 5: Final Review", [
    "User is presented with a preview of their entire complaint.",
    "Can review answers and check uploaded files.",
    "Upon confirmation, the Automated PDF is generated.",
    "Complaint is officially logged in the database."
])

# Slide 23
add_content_slide(prs, "22. Authentication Flow - Overview", [
    "Security is paramount in cybercrime reporting.",
    "Platform offers multiple layers of authentication.",
    "Ensures the identity of the reporter is known and verifiable.",
    "Prevents spam and fake reporting."
])

# Slide 24
add_content_slide(prs, "23. Auth Flow: Traditional Login", [
    "Users can register with Email, Name, and Phone Number.",
    "Passwords are securely hashed using Werkzeug.",
    "Session management ensures secure access to the dashboard.",
    "OTP verification can be enforced for added security."
])

# Slide 25
add_content_slide(prs, "24. Auth Flow: SSO & DigiLocker", [
    "Google SSO integration for 1-click secure login.",
    "DigiLocker integration to verify official government documents.",
    "Elevates the trust level of the user's account.",
    "Streamlines the onboarding process for citizens."
])

# Slide 26
add_content_slide(prs, "25. Conclusion & Future Scope", [
    "CyberCrimeAI modernizes the reporting infrastructure.",
    "Future Scope: Direct integration with police precinct databases.",
    "Future Scope: Multi-lingual AI support for rural areas.",
    "Thank you for reviewing the CyberCrimeAI platform!"
])


# Save the presentation
prs.save('CyberCrimeAI_Presentation_25Pages.pptx')
print("Successfully generated CyberCrimeAI_Presentation_25Pages.pptx")
